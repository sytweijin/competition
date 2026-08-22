"""C3：PDF/Word/PPT/Excel/图片/音频全格式解析链路实测。"""

from __future__ import annotations

import io
import wave

import pytest
from PIL import Image

from app.file_analysis import extract_text


@pytest.fixture(autouse=True)
def _disable_real_media_calls(monkeypatch):
    import app.services.media_analysis as media

    monkeypatch.setattr(media, "APP_VISION_API_KEY", "")
    monkeypatch.setattr(media, "APP_VISION_MODEL", "")
    monkeypatch.setattr(media, "APP_ASR_API_KEY", "")
    monkeypatch.setattr(media, "APP_ASR_MODEL", "")
    monkeypatch.setattr(media, "MAP_REALTIME_API_KEY", "")
    monkeypatch.setattr(media, "ASCEND_OMNI_WS_URL", "")


def _png_bytes() -> bytes:
    buffer = io.BytesIO()
    Image.new("RGB", (4, 4), (255, 255, 255)).save(buffer, format="PNG")
    return buffer.getvalue()


def _jpg_bytes() -> bytes:
    buffer = io.BytesIO()
    Image.new("RGB", (4, 4), (0, 128, 255)).save(buffer, format="JPEG")
    return buffer.getvalue()


def _wav_bytes() -> bytes:
    buffer = io.BytesIO()
    with wave.open(buffer, "wb") as wav:
        wav.setnchannels(1)
        wav.setsampwidth(2)
        wav.setframerate(8000)
        wav.writeframes(b"\x00\x00" * 800)
    return buffer.getvalue()


def _pdf_bytes() -> bytes:
    from reportlab.pdfgen import canvas

    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer)
    pdf.drawString(72, 720, "Final report for campus survey")
    pdf.showPage()
    pdf.save()
    return buffer.getvalue()


def _docx_bytes() -> bytes:
    from docx import Document

    buffer = io.BytesIO()
    doc = Document()
    doc.add_paragraph("完成调研问卷和数据分析")
    doc.save(buffer)
    return buffer.getvalue()


def _xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    buffer = io.BytesIO()
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "调研问卷"
    ws["B1"] = 20
    wb.save(buffer)
    return buffer.getvalue()


def _pptx_bytes() -> bytes:
    from pptx import Presentation

    buffer = io.BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "答辩 PPT 初稿"
    prs.save(buffer)
    return buffer.getvalue()


DOCUMENT_CASES = [
    ("需求说明.pdf", _pdf_bytes, "Final report"),
    ("需求说明.docx", _docx_bytes, "调研问卷"),
    ("需求说明.xlsx", _xlsx_bytes, "调研问卷"),
    ("需求说明.pptx", _pptx_bytes, "答辩 PPT"),
]

IMAGE_CASES = [
    ("photo.png", _png_bytes),
    ("photo.jpg", _jpg_bytes),
    ("photo.jpeg", _png_bytes),
    ("photo.webp", _png_bytes),
]

AUDIO_CASES = [
    ("voice.mp3", lambda: b"ID3fake-audio"),
    ("voice.wav", _wav_bytes),
    ("voice.m4a", lambda: b"fake-m4a"),
    ("voice.webm", lambda: b"fake-webm"),
]


@pytest.mark.parametrize("filename,make_content,expected", DOCUMENT_CASES)
def test_document_formats_extract_text(filename, make_content, expected):
    text = extract_text(filename, make_content())
    assert expected in text


@pytest.mark.parametrize("filename,make_content", IMAGE_CASES)
def test_image_formats_fallback_without_model(filename, make_content):
    text = extract_text(filename, make_content())
    assert "图片文件" in text
    assert "未配置视觉模型" in text


@pytest.mark.parametrize("filename,make_content", AUDIO_CASES)
def test_audio_formats_fallback_without_model(filename, make_content):
    text = extract_text(filename, make_content())
    assert "音频文件" in text
    assert "未配置语音转写模型" in text


def test_all_audio_formats_forward_to_asr_with_model(monkeypatch):
    import app.services.media_analysis as media

    class FakeTranscriptions:
        def create(self, **kwargs):
            return type("Resp", (), {"text": "会议录音转写结果"})

    class FakeAudio:
        def __init__(self):
            self.transcriptions = FakeTranscriptions()

    class FakeClient:
        def __init__(self):
            self.audio = FakeAudio()

    monkeypatch.setattr(media, "APP_ASR_API_KEY", "key")
    monkeypatch.setattr(media, "APP_ASR_MODEL", "whisper-1")
    monkeypatch.setattr(media, "APP_ASR_TRANSCRIPTION_MODE", "native")
    monkeypatch.setattr(
        media, "_client", lambda api_key, base_url: FakeClient())
    for filename, make_content in AUDIO_CASES:
        text = extract_text(filename, make_content())
        assert "音频转写" in text
        assert "会议录音转写结果" in text


def test_all_image_formats_forward_to_vision_with_model(monkeypatch):
    import app.services.media_analysis as media

    class FakeMessage:
        content = "图片中的文字：宣传海报"

    class FakeChoice:
        message = FakeMessage()

    class FakeCompletions:
        def create(self, **kwargs):
            return type("Resp", (), {"choices": [FakeChoice()]})

    class FakeClient:
        def __init__(self):
            self.chat = type("Chat", (), {"completions": FakeCompletions()})()

    monkeypatch.setattr(media, "APP_VISION_API_KEY", "key")
    monkeypatch.setattr(media, "APP_VISION_MODEL", "gpt-4o")
    monkeypatch.setattr(
        media, "_client", lambda api_key, base_url: FakeClient())
    for filename, make_content in IMAGE_CASES:
        text = extract_text(filename, make_content())
        assert "图片理解" in text
        assert "宣传海报" in text
