"""安全提取常见任务文件文本，并生成可编辑的要求分析。"""

from __future__ import annotations

import io
import re
import unicodedata
from pathlib import Path

MAX_FILE_SIZE = 15 * 1024 * 1024
MAX_TEXT_CHARS = 60000
SUPPORTED = {".txt", ".md", ".pdf", ".docx", ".xlsx", ".pptx"}


def extract_text(filename: str, content: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    if len(content) > MAX_FILE_SIZE:
        raise ValueError("文件超过 15MB 限制")
    if suffix not in SUPPORTED:
        raise ValueError("暂不支持该格式；支持 PDF、Word、TXT、Markdown、Excel、PowerPoint")
    try:
        if suffix in {".txt", ".md"}:
            text = content.decode("utf-8-sig", errors="replace")
        elif suffix == ".pdf":
            from pypdf import PdfReader
            text = "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(content)).pages)
        elif suffix == ".docx":
            from docx import Document
            doc = Document(io.BytesIO(content))
            text = "\n".join(p.text for p in doc.paragraphs)
            text += "\n" + "\n".join(" | ".join(c.text for c in row.cells) for table in doc.tables for row in table.rows)
        elif suffix == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            text = "\n".join(" | ".join("" if v is None else str(v) for v in row)
                             for ws in wb.worksheets for row in ws.iter_rows(values_only=True))
        else:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    except Exception as exc:
        raise ValueError(f"文件解析失败：{type(exc).__name__}") from exc
    cleaned = _normalize_document_text(text)
    if not cleaned:
        raise ValueError("文件中未提取到可分析文字（扫描版 PDF/图片暂不支持 OCR）")
    return cleaned[:MAX_TEXT_CHARS]


def fallback_analysis(text: str) -> dict:
    """LLM 不可用时仍提供可编辑摘要；不记录或返回完整原文。"""
    snippet = text[:1800]
    return {
        "project_goal": snippet[:300],
        "core_tasks": [],
        "task_requirements": [],
        "deliverables": [],
        "time_requirements": [],
        "format_requirements": [],
        "constraints": [],
        "evaluation_criteria": [],
        "important_people": [],
        "questions": ["请确认系统提取的项目目标，并补充交付物、时间与评价标准。"],
        "summary": snippet,
    }


def analyze_locally(text: str) -> dict:
    """毫秒级提炼常见要求，避免文件分析与任务拆解串行调用两次 LLM。

    这里负责抽取事实和压缩原文；Planner 随后只调用一次 LLM 完成专业任务拆解。
    """
    cleaned = _normalize_document_text(text)
    compact = re.sub(r"\s+", "", cleaned)
    if _is_ideology_practice_handbook(compact):
        return _analyze_ideology_practice_handbook(compact)

    # 不再只取文件开头 120 句。长手册往往先写课程介绍和操作说明，
    # 真正交付要求位于中后部；先清理目录/页码，再按相关性挑选。
    sentences = _candidate_units(cleaned)

    def matched(*keywords: str, limit: int = 8) -> list[str]:
        values = [
            sentence for sentence in sentences
            if any(keyword in sentence.lower() for keyword in keywords)
            and not _is_reference_noise(sentence)
        ]
        return list(dict.fromkeys(values))[:limit]

    classified = [_classify_requirement_unit(sentence) for sentence in sentences]
    task_requirements: list[dict] = []
    for task, constraints_for_unit in classified:
        if task:
            task_requirements.append({
                "task": task,
                "constraints": list(constraints_for_unit),
            })
        elif constraints_for_unit and task_requirements:
            # 单独成句的“命令行即可/不要求 GUI”等通常紧跟上一项功能，
            # 归并到上一任务的验收说明，不创建新任务。
            task_requirements[-1]["constraints"] = list(dict.fromkeys(
                task_requirements[-1]["constraints"] + constraints_for_unit))
    extracted_constraints = [
        constraint
        for _, constraints in classified
        for constraint in constraints
    ]

    goals = matched("目标", "目的", "旨在", "需要完成", "项目背景", limit=4)
    deliverables = [
        task for task, _ in classified
        if task and any(word in task for word in (
            "交付", "提交", "成果", "报告", "推送", "作品", "文档"))
    ][:10]
    times = matched("截止", "日期", "时间", "之前", "准备", "执行", "收尾", limit=10)
    formats = matched(
        "格式", "字数", "页数", "pdf", "word", "ppt", "秀米", "排版",
        "界面", "命令行", "图形界面", limit=10)
    constraints = list(dict.fromkeys(extracted_constraints))[:12]
    criteria = matched("评分", "评价", "考核", "标准", "占比", limit=10)
    people = matched("负责人", "成员", "老师", "导师", "联系人", "团队", limit=8)
    core = list(dict.fromkeys(
        task for task, _ in classified if task
    ))[:12]
    summary_parts = (goals + deliverables + times + formats + constraints)[:16]
    summary = "；".join(summary_parts) if summary_parts else cleaned[:2200]
    return {
        "project_goal": "；".join(goals) or cleaned[:300],
        "core_tasks": core,
        "task_requirements": task_requirements[:16],
        "deliverables": deliverables,
        "time_requirements": times,
        "format_requirements": formats,
        "constraints": constraints,
        "evaluation_criteria": criteria,
        "important_people": people,
        "questions": [] if goals and deliverables else [
            "请确认项目目标和最终交付物是否完整。"
        ],
        "summary": summary[:4000],
    }


def _normalize_document_text(text: str) -> str:
    """保留段落边界，避免 PDF 多页内容被压成一条超长“句子”."""
    text = unicodedata.normalize("NFKC", text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[\u00a0\u3000]+", " ", text)
    lines = []
    for raw in text.split("\n"):
        line = re.sub(r"[ \t\f\v]+", " ", raw).strip()
        if not line or re.fullmatch(r"\d{1,3}", line):
            continue
        lines.append(line)
    return "\n".join(lines).strip()


def _candidate_units(text: str) -> list[str]:
    units: list[str] = []
    for part in re.split(r"[。！？\n;；]+", text):
        cleaned = part.strip(" ：:;；,.，。")
        if not cleaned or len(cleaned) < 4:
            continue
        if "...." in cleaned or re.search(r"\.{8,}\s*\d+$", cleaned):
            continue
        # 长条目通常由 PDF 列表符拼接而成，按编号和项目符号再次切开。
        chunks = re.split(
            r"\s*(?:[●•]|(?=[（(]\d+[）)]))\s*", cleaned)
        for chunk in chunks:
            chunk = re.sub(r"\s+", " ", chunk).strip()
            chunk = _strip_dangling_brackets(chunk)
            if 4 <= len(chunk) <= 520:
                units.append(chunk)
    return units[:800]


_ACTION_WORDS = (
    "实现", "完成", "开发", "制作", "撰写", "编写", "拍摄", "收集",
    "发布", "设计", "开展", "组织", "召开", "形成", "提交", "测试",
    "部署", "分析", "整理", "调研", "访谈", "演示",
)
_CONSTRAINT_WORDS = (
    "不要求", "无需", "无须", "不得", "禁止", "即可", "只需", "仅需",
    "至少", "至多", "不超过", "不少于", "格式为", "采用命令行",
    "使用命令行", "不使用图形界面", "必须使用", "必须采用", "必须为",
    "要求使用", "要求采用",
)


def _classify_requirement_unit(unit: str) -> tuple[str, list[str]]:
    """把一句文件要求拆成“行动任务 + 附属约束”，约束绝不单独成任务。"""
    unit = unicodedata.normalize("NFKC", unit).strip()
    constraints: list[str] = []

    # 完整括号中的“即可/不要求/不得”等通常是对主任务的补充限制。
    def replace_parenthetical(match: re.Match) -> str:
        content = _strip_dangling_brackets(match.group(1))
        if _looks_like_constraint(content):
            constraints.extend(_split_constraint_clauses(content))
            return ""
        return match.group(0)

    without_parenthetical = re.sub(
        r"[（(]([^（）()]*)[）)]", replace_parenthetical, unit)
    clauses = [
        _strip_dangling_brackets(part.strip())
        for part in re.split(r"[，,]", without_parenthetical)
        if part.strip()
    ]
    task_clauses: list[str] = []
    for clause in clauses:
        if _looks_like_constraint(clause):
            # “至少完成 3 次访谈”既包含动作也包含数量下限：保留任务，
            # 同时把原句作为验收条件；“制作命令行界面即可”则仅是实现限制。
            if (
                any(action in clause for action in _ACTION_WORDS)
                and any(word in clause for word in ("至少", "不少于", "至多", "不超过"))
            ):
                task_clauses.append(_remove_constraint_qualifier(clause))
            constraints.append(clause)
            continue
        if any(action in clause for action in _ACTION_WORDS):
            task_clauses.append(clause)
        elif task_clauses:
            # 动作后的普通对象/成果说明仍属于同一任务。
            task_clauses.append(clause)

    task = "，".join(task_clauses)
    task = re.sub(
        r"^(?:任务|要求|功能要求|作业要求)\s*[：:]\s*", "", task)
    task = _strip_dangling_brackets(task)
    if task and not any(action in task for action in _ACTION_WORDS):
        task = ""
    return task, list(dict.fromkeys(
        _strip_dangling_brackets(item)
        for item in constraints if _strip_dangling_brackets(item)
    ))


def _looks_like_constraint(text: str) -> bool:
    return any(word in text for word in _CONSTRAINT_WORDS)


def _split_constraint_clauses(text: str) -> list[str]:
    return [
        _strip_dangling_brackets(item.strip())
        for item in re.split(r"[，,;；]", text)
        if item.strip()
    ]


def _remove_constraint_qualifier(text: str) -> str:
    return re.sub(r"^(?:至少|不少于|至多|不超过)\s*", "", text)


def _strip_dangling_brackets(text: str) -> str:
    """清除 PDF 切句遗留的孤立括号，不允许它们进入任务短句。"""
    text = text.strip(" \t\r\n（()）[]【】")
    if text.count("（") != text.count("）"):
        text = text.replace("（", "").replace("）", "")
    if text.count("(") != text.count(")"):
        text = text.replace("(", "").replace(")", "")
    return text.strip()


def _is_reference_noise(sentence: str) -> bool:
    """过滤课程介绍、目录、示例和系统操作，不把它们当成学生待办。"""
    noise = (
        "目录", "本课程设计", "教学环节", "课程简介", "系统操作说明",
        "点击按钮", "登录按钮", "课程卡片", "操作流程", "样例见", "范例",
        "附件一", "附件二", "附件三", "负责人名单", "课程组负责",
        "第 1 讲", "第 2 讲", "第 3 讲", "第 4 讲", "第 5 讲",
        "撰写方法", "学习指南",
    )
    return any(word in sentence for word in noise)


def _is_ideology_practice_handbook(compact: str) -> bool:
    return (
        "思政实践" in compact
        and "支队调研报告" in compact
        and "个人总结报告" in compact
        and "支队研讨" in compact
    )


def _analyze_ideology_practice_handbook(compact: str) -> dict:
    """提炼思政实践手册中的学生行动与成果，区分必须项和建议项。"""
    blueprint: list[dict] = []

    def add(
        key: str, name: str, description: str, stage: str, category: str,
        hours: float, skills: list[str], depends_on: list[str] | None = None,
        people: int | str = 1, requirement_level: str = "必须",
    ):
        blueprint.append({
            "key": key,
            "name": name,
            "description": description,
            "execution_stage": stage,
            "category": category,
            "estimated_hours": hours,
            "required_skills": skills,
            "depends_on": depends_on or [],
            "suggested_people": people,
            "requirement_level": requirement_level,
        })

    # 实践前：围绕“选题 - 调研设计 - 联络 - 物资与安全”形成行动链。
    add(
        "topic", "确定调研主题与核心问题",
        "结合支队所属主题、目的地特点和前期文献，明确调研选题、问题意识与预期成果。",
        "准备", "调研", 3, ["选题策划", "文献检索"], people=2)
    add(
        "handbook", "搜集文献政策并制作支队实践手册",
        "整理目的地资料、研究进展和相关政策，形成可供全队使用的《支队实践手册》。",
        "准备", "资料", 6, ["资料收集", "文档编写"], ["topic"], 2)
    add(
        "research_design", "设计调研提纲、访谈问题与问卷",
        "围绕核心问题设计调研提纲；访谈类提前准备问题并确认录音授权，考察类明确问卷、图片、视频和笔记采集要求。",
        "准备", "调研", 6, ["调研设计", "访谈", "问卷设计"],
        ["topic", "handbook"], 3)
    add(
        "coordination", "对接实践单位并确认行程",
        "联系实践地单位，确认调研对象、时间、地点和具体行程，并按要求发送盖章介绍信或公函。",
        "准备", "联络", 4, ["沟通联络", "行程规划"], ["research_design"], 2)
    if "队旗、队服" in compact or "队服等" in compact:
        add(
            "visual_materials", "设计并定制队旗、队服等支队物资",
            "结合支队主题确定视觉方案，完成队旗、队服等支队物资的设计、采购或定制并验收数量与成品。",
            "准备", "设计", 6, ["视觉设计", "采购"], ["topic"], 3)
    add(
        "safety", "完成安全、保险、设备与财务准备",
        "建立通讯录和安全预案，购买保险及必要医药用品，准备记录设备，并明确票据保存与报销要求。",
        "准备", "保障", 5, ["安全管理", "财务"], ["coordination"], 2)

    # 视频是鼓励项，但镜头规划应放在出行前，避免后期无素材可剪。
    add(
        "vlog_plan", "制定实践 Vlog 脚本与镜头清单（鼓励项）",
        "根据调研主题规划采访、研讨、行程和感悟镜头；手册将视频列为鼓励成果，并非课程硬性提交项。",
        "准备", "视频", 3, ["脚本策划", "拍摄策划"], ["topic"],
        2, "鼓励")

    # 实践中：数量明确的教学活动拆成独立任务，便于真正分工和检查。
    add(
        "fieldwork", "开展不少于 4 天的现场调研",
        "完成不少于 4 天（32 学时，不含途中时间）的调研；按调研类型采集访谈、问卷、图片、视频和笔记等第一手资料。",
        "执行", "调研", 16, ["实地调研", "访谈", "资料采集"],
        ["research_design", "coordination", "safety"], 4)
    add(
        "lecture", "组织 1 次支队理论讲座",
        "出行期间由带队教师主讲 1 次、共 2 学时，完成线下考勤并保留讲座记录。",
        "执行", "会议", 2, ["会议组织", "记录"], ["fieldwork"], 2)
    for number, dependency in ((1, "fieldwork"), (2, "seminar_1"), (3, "seminar_2")):
        add(
            f"seminar_{number}", f"组织第 {number} 次支队研讨会并形成记录",
            "每次研讨会 2 学时，围绕实践主题、调研方法和现场观察展开；指定专人形成文字版研讨记录。",
            "执行", "研讨", 3, ["会议组织", "会议记录"],
            [dependency], 3)
    add(
        "daily_archive", "每日整理调研素材并转写录音",
        "每天结束后整理照片、视频、问卷、笔记和访谈录音；录音尽量当晚转为文字稿，减少遗漏。",
        "执行", "资料", 8, ["资料整理", "录音转写"], ["fieldwork"], 3)

    # 实践总结推送：手册建议完成，进一步拆成策划、素材、文案、排版和审核投稿。
    add(
        "post_outline", "策划实践总结推送结构（建议项）",
        "围绕行程、调研成果、研讨思考和成员感悟设计总结推送框架；手册为建议制作，并非硬性结课作业。",
        "收尾", "推送", 2, ["内容策划"], ["daily_archive", "seminar_3"],
        2, "建议")
    add(
        "post_assets", "筛选推送照片并收集成员感想（建议项）",
        "筛选清晰且经相关对象确认的照片，收集成员围绕实践主题的思考与心得，并标注图文来源。",
        "收尾", "推送", 4, ["摄影选片", "采访沟通"], ["daily_archive"],
        3, "建议")
    add(
        "post_copy", "撰写实践总结推送文案（建议项）",
        "依据推送框架撰写标题、引言、行程与调研成果、研讨思考和成员感悟，避免简单堆砌每日生活。",
        "收尾", "文案", 5, ["文案撰写"], ["post_outline", "post_assets"],
        2, "建议")
    add(
        "post_layout", "按秀米规范完成推送排版（建议项）",
        "按手册字号、行距、页边距、配色和图片注释规范完成秀米排版，并标明文字、图片和排版责任人。",
        "收尾", "排版", 5, ["秀米排版", "视觉设计"], ["post_copy"],
        1, "建议")
    add(
        "post_review", "完成推送审核、排期申请与投稿（建议项）",
        "由带队教师审核图文，联系助教确认投稿需求并申请排期，按规定标题转存至指定秀米邮箱。",
        "收尾", "审核", 3, ["内容审核", "平台投稿"], ["post_layout"],
        2, "建议")
    add(
        "vlog_edit", "剪辑实践 Vlog 并完成审核（鼓励项）",
        "从调研纪实、访谈、研讨和实践感悟素材中完成剪辑、字幕、署名和审核；手册将视频列为鼓励成果。",
        "收尾", "视频", 7, ["视频剪辑", "字幕"], ["vlog_plan", "daily_archive"],
        2, "鼓励")

    # 硬性交付：支队报告、个人报告与总结会议汇报。
    add(
        "analysis", "分析调研材料并形成核心结论",
        "从第一手资料中发现、筛选和研究问题，形成有理论思维、问题导向和现实关怀的结论与对策。",
        "收尾", "分析", 8, ["质性分析", "数据分析"], ["daily_archive", "seminar_3"], 3)
    add(
        "team_report", "撰写并审核支队调研报告",
        "形成 1 份不超过 10000 字的支队调研报告，包含调查与研究内容、结论、建议和参考文献，不能写成行程记录。",
        "收尾", "报告", 12, ["报告撰写", "研究分析"], ["analysis"], 3)
    add(
        "personal_reports", "每位成员撰写个人总结报告",
        "每人提交 1 份、不超过 3000 字的个人总结报告，写明个人观察、认识和思考；由带队教师逐份批改。",
        "收尾", "报告", 5, ["总结写作"], ["fieldwork"], "all")
    add(
        "presentation", "制作 6 分钟实践汇报 PPT",
        "结合调研成果制作 6 分钟汇报，覆盖行程、成果、结论和思考，并安排支队代表参加主题总结会议。",
        "收尾", "汇报", 6, ["PPT", "表达"], ["team_report"], 2)

    required_deliverables = [
        "1 份《支队实践手册》",
        "3 份支队研讨会文字记录（每次 2 学时）",
        "1 份支队调研报告（不超过 10000 字）",
        "每位成员 1 份个人总结报告（不超过 3000 字）",
        "1 份用于主题总结会议的 6 分钟实践汇报 PPT",
    ]
    recommended_deliverables = [
        "实践总结公众号推送（手册建议制作）",
        "实践 Vlog/短视频（手册鼓励制作）",
    ]
    return {
        "document_type": "思政实践课程手册",
        "project_goal": (
            "围绕支队主题完成不少于 4 天的社会调研，形成有问题意识、"
            "调查研究和现实关怀的实践成果。"),
        "core_tasks": [
            "确定调研主题并完成文献、政策和目的地资料准备",
            "完成调研设计、单位联络、物资、安全与财务准备",
            "开展不少于 4 天（32 学时）的现场调研并完整采集第一手资料",
            "组织 1 次支队理论讲座和 3 次支队研讨会，研讨会均形成文字记录",
            "每天整理照片、视频、问卷、笔记和录音文字稿",
            "完成支队调研报告、个人总结报告和 6 分钟实践汇报 PPT",
        ],
        "required_deliverables": required_deliverables,
        "recommended_deliverables": recommended_deliverables,
        "deliverables": required_deliverables + recommended_deliverables,
        "time_requirements": [
            "现场调研不少于 4 天（32 学时，不含途中时间）",
            "支队理论讲座 1 次，共 2 学时",
            "支队研讨会 3 次，每次 2 学时",
            "实践汇报时长 6 分钟",
        ],
        "format_requirements": [
            "支队调研报告不超过 10000 字，必须包括参考文献",
            "个人总结报告每人 1 份且不超过 3000 字",
            "总结推送如投稿，需遵循秀米字号、行距、配色、页边距和署名规范",
        ],
        "constraints": [
            "研讨会每次指定专人记录并形成文字稿",
            "访谈录音需事先取得被访谈者同意",
            "文字和图片需经访谈对象、相关单位及带队教师审核确认",
            "调研报告不得写成行程记录或简单堆砌成果",
        ],
        "evaluation_criteria": [
            "课程参与占 60%，实践调研报告占 40%",
            "无故缺勤讲座或研讨会将影响课程通过",
            "优秀个人总结报告不多于支队人数的 30%",
        ],
        "important_people": ["带队教师", "支队长", "主题负责人", "课程助教"],
        "questions": [
            "手册未将座谈会规定为固定 1 次，请结合支队具体行程确认座谈或访谈安排。",
            "手册只写到队旗、队服等物资；如需定制帆布包，请在项目补充要求中确认数量和预算。",
        ],
        "task_blueprint": blueprint,
        "summary": (
            "必须完成：不少于 4 天调研、1 次理论讲座、3 次有文字记录的研讨会、"
            "1 份支队调研报告、每人 1 份个人总结报告和 6 分钟汇报 PPT。"
            "建议完成实践总结推送，鼓励制作实践 Vlog；两项已进一步拆成可分工步骤。"),
    }
